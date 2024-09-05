#!/usr/bin/env python
import argparse
import enum
import glob
import itertools
import os
import re
import sys
import xml.etree.ElementTree as ET
from typing import List, Callable, Optional
from zipfile import ZipFile

import bs4
import cmd2
import pptx
from lxml import etree

import pptxutil
from pptxutil import PPTXHelper

TEMPLATE_PATH = 'templates'


class WorkItem(enum.Enum):
    presentation = enum.auto()
    template = enum.auto()


def workon_template(method: Callable):
    """
    Decorator for methods that work on templates
    :param method:
    :return:
    """

    def decorator(ppt_app: 'PPTApp', *args, **kwargs):
        if ppt_app._work_on != WorkItem.template:
            ppt_app.perror('Need to select a template first')
            return
        return method(ppt_app, *args, **kwargs)

    return cmd2.with_category('PPT')(decorator)


class PPTApp(cmd2.Cmd):

    def __init__(self):
        super().__init__()
        self._work_on = WorkItem.template
        self._work_item = self._templates()[2]
        self._set_prompt()

    @property
    def template_path(self):
        return os.path.join(TEMPLATE_PATH, f'{self._work_item}.pptx')

    def _set_prompt(self):
        self.prompt = f'({self._work_item}) '

    def postcmd(self, stop: bool, line: str) -> bool:
        """Hook method executed just after a command dispatch is finished.
        :param stop: if True, the command has indicated the application should exit
        :param line: the command line text for this command
        :return: if this is True, the application will exit after this command and the postloop() will run
        """
        self._set_prompt()
        return stop

    def _templates(self) -> List[str]:
        """
        Return a list of template names
        :return:
        """
        templates = glob.glob(f'{TEMPLATE_PATH}/*.pptx')
        templates = list(map(lambda name: os.path.splitext(os.path.basename(name))[0], templates))
        templates.sort()
        return templates

    def _template_complete(self, text: str, line: str, begidx: int, endidx: int):
        """
        Completer method for template names
        :param text:
        :param line:
        :param begidx:
        :param endidx:
        :return:
        """
        templates = self._templates()
        return [t for t in templates if t.startswith(text)]

    parser = argparse.ArgumentParser()
    parser.add_argument('template', help='template name', completer=_template_complete)

    @cmd2.with_category('PPT')
    @cmd2.with_argparser(parser)
    def do_template(self, arg: argparse.Namespace):
        """
        Select template to work on
        """
        template = arg.template
        if template not in self._templates():
            self.perror(f'Template "{template}" not found. Available: {", ".join(t for t in self._templates())}')
            return
        self._work_item = template
        self._work_on = WorkItem.template

    def _ppt_presentation(self) -> Optional[pptx.Presentation]:
        """
        Get Presentation object for selected template or presentation
        :return:
        """
        if self._work_on == WorkItem.template:
            return pptx.Presentation(self.template_path)
        else:
            return pptx.Presentation(self._work_item)

    def _layout_complete(self, text: str, line: str, begidx: int, endidx: int):
        """
        Completer method for layout names within selected presentation or template
        :param text:
        :param line:
        :param begidx:
        :param endidx:
        :return:
        """
        ppt = self._ppt_presentation()
        layouts = [l.name for l in ppt.slide_layouts
                   if l.name.startswith(text)]
        return layouts

    parser = argparse.ArgumentParser()
    parser.add_argument('layout', help='layout name', nargs='?', default=None, completer=_layout_complete)

    @cmd2.with_category('PPT')
    @cmd2.with_argparser(parser)
    def do_layout(self, arg: argparse.Namespace):
        """
        List layouts in template
        :return:
        """
        ppt = self._ppt_presentation()
        layout_name = arg.layout

        for index, layout in enumerate(ppt.slide_layouts):
            if layout_name and layout_name != layout.name:
                continue
            self.poutput(f'Layout: {layout.name}')
            # Go through all the placeholders and identify them by index and type
            shapes = [s for s in layout.shapes]
            phs = [p for p in layout.placeholders]
            self.show_place_holders(layout)

    parser = argparse.ArgumentParser()
    parser.add_argument('layout', help='layout name', completer=_layout_complete)

    @cmd2.with_category('PPT')
    @cmd2.with_argparser(parser)
    def do_sort_layout(self, arg: argparse.Namespace):
        """
        Sort placeholders in layout
        :return:
        """
        ppt = self._ppt_presentation()
        layout_name = arg.layout
        layout = next((l for l in ppt.slide_layouts if l.name == layout_name), None)
        if layout is None:
            self.perror(f'Layout "{layout}" not found')
            return

        def pretty(xml_str: str) -> str:
            import xml.dom.minidom

            dom = xml.dom.minidom.parseString(xml_str)  # or xml.dom.minidom.parseString(xml_string)
            pretty_xml_as_string = dom.toprettyxml()
            return pretty_xml_as_string

        layout_xml = pretty(ET.tostring(layout.element))
        placeholder_xml = pretty(ET.tostring(layout.placeholders.element))
        self.poutput(placeholder_xml)

    def show_place_holders(self, layout):
        shapes = [(s, i) for i, s in enumerate(layout.placeholders)]
        if not shapes:
            self.poutput('  no shapes')
            return
        shapes.sort(key=lambda tpl: (tpl[0].top, tpl[0].left))
        max_name = max(len(shp.name) for shp, _ in shapes)
        max_text = max(len(shp.text) for shp, _ in shapes)
        max_index = max(len(f'{idx}') for _, idx in shapes)
        for shape, shp_index in shapes:
            self.poutput(f'  shape {shp_index:{max_index}}: name {shape.name:{max_name}} text {shape.text:{max_text}}')

    @cmd2.with_category('PPT')
    def do_master(self, arg: argparse.Namespace):
        """
        Show info about master slide
        """
        ppt = self._ppt_presentation()
        master = ppt.slide_master
        self.show_place_holders(master)

    @cmd2.with_category('PPT')
    def do_zipinfo(self, arg: argparse.Namespace):
        pptx = pptxutil.PPTXHelper(file=self.template_path)
        themes_and_color_specs = pptx.themes_color_specs()
        for name, color_spec in themes_and_color_specs.values():
            if name == 'Office':
                continue
            color_len = max(len(c) for c in color_spec)
            self.poutput(f'Theme: {name}')
            for c_name, c_spec in color_spec.items():
                self.poutput(f'  {c_name:{color_len}}: {c_spec}')
        files = pptx.namelist()
        file_info = pptx.infolist()
        foo = 1

    def _pptx_completer(self, text: str, line: str, begidx: int, endidx: int):
        """
        Completer for PPTX files in the current directory
        :param text:
        :param line:
        :param begidx:
        :param endidx:
        :return:
        """

        def pptx_filter(path: str):
            """
            Only allow paths and PPTX files
            :param path:
            :return:
            """
            return os.path.isdir(path) or os.path.isfile(path) and os.path.splitext(path)[1].lower() == '.pptx'

        options = cmd2.Cmd.path_complete(self, text=text, line=line, begidx=begidx, endidx=endidx,
                                         path_filter=pptx_filter)
        return options

    parser = argparse.ArgumentParser()
    parser.add_argument('presentation', help='presentation to work on', completer=_pptx_completer)

    @cmd2.with_category('PPT')
    @cmd2.with_argparser(parser)
    def do_update_colors(self, arg: argparse.Namespace):
        """
        Change colors in ppt from scheme to RGB
        :param arg:
        :return:
        """
        # open PPT
        # get scheme colors
        # for each slide
        # replace color specs that reference a scheme color to an RGB color
        ppt_path = arg.presentation
        if not os.path.isfile(ppt_path):
            self.perror(f'File "{ppt_path}" not found.')
            return
        if os.path.splitext(ppt_path)[1].lower() != '.pptx':
            self.perror(f'File "{ppt_path}" is not a PPTX file.')
            return
        pptx = pptxutil.PPTXHelper(file=ppt_path)
        themes_and_color_specs = pptx.themes_color_specs()
        # get a sorted list of all color spec tags from all themes
        color_tags = sorted(
            set(itertools.chain.from_iterable((c_specs.keys() for name, c_specs in themes_and_color_specs.values()))))

        # print headline with all theme names
        theme_names = [name for name, _ in themes_and_color_specs.values()]
        name_len = max(map(len, theme_names))
        name_len = max(7, name_len)
        theme_fnames = [n.split('/')[-1].split('.')[0] for n in themes_and_color_specs]
        print(f'{"":8}  {" ".join(f"{n:{name_len}}" for n in theme_fnames)}')
        print(f'{"":8}  {" ".join(f"{n:{name_len}}" for n in theme_names)}')

        # now for each tag print a row with the value from each theme
        for color_tag in color_tags:
            spec = ' '.join((f'{s.get(color_tag, ""):{name_len}}' for name, s in themes_and_color_specs.values()))
            print(f'{color_tag:8}: {spec}')

        slides = pptx.slide_file_info_list()
        slides.sort(key=lambda x: int(re.search(r'slide(\d+)\.xml', x.filename).group(1)))
        slide_info = []
        for slide in slides:
            with pptx.open(slide) as slide_file:
                slide_data = slide_file.read()
            soup = bs4.BeautifulSoup(slide_data, 'xml')
            clrMapOvr = soup.find('p:clrMapOvr')
            srgb_clr_elements = soup.find_all('a:srgbClr')
            scheme_clr_elements = soup.find_all('a:schemeClr')
            slide_info.append((slide.filename.split('/')[-1].split('.')[0],
                               {'clr_map_ovr': clrMapOvr,
                                'srgb_clr_elements': srgb_clr_elements,
                                'scheme_clr_elements': scheme_clr_elements}))

        slides_with_override = [(slide, slide_data) for slide, slide_data in slide_info
                                if slide_data['clr_map_ovr'].find('a:overrideClrMapping')]
        slides_wo_override = [(slide, slide_data) for slide, slide_data in slide_info
                              if slide_data['clr_map_ovr'].find('a:masterClrMapping')]
        srgb_parents = sorted(
            set(itertools.chain.from_iterable(
                ((e.parent.name for e in sd['srgb_clr_elements']) for _, sd in slide_info))))
        scheme_clr_parents = sorted(set(itertools.chain.from_iterable((('{}:{}:{}'.format(e.parent.parent.parent.name,
                                                                                          e.parent.parent.name,
                                                                                          e.parent.name) for e in
                                                                        sd['scheme_clr_elements']) for _, sd in
                                                                       slide_info))))

        files = pptx.namelist()
        file_info = pptx.infolist()
        presentation_xml_rels = next((f for f in file_info if f.filename.endswith('presentation.xml.rels')))
        with pptx.open(presentation_xml_rels) as presentation_xml_rels_file:
            rel_data = presentation_xml_rels_file.read()
        soup = bs4.BeautifulSoup(rel_data, 'xml')
        relationships = [r for r in soup.Relationships.children]
        relationship_types = sorted(set((r['Type'].split('/')[-1] for r in relationships)))
        theme_relationships = [r for r in relationships if r['Type'].endswith('theme')]

        paths = set(('/'.join(f.filename.split('/')[:-1]) for f in file_info))
        themes = [f for f in file_info if f.filename.startswith('ppt/theme')]
        for theme in themes:
            with pptx.open(theme) as theme_file:
                theme_data = theme_file.read()
            tree = ET.fromstring(theme_data)
            m = re.search(r'\{(.+)\}theme', tree.tag)
            ns = {'ns': m.group(1)}
            clr_scheme = next(tree.iterfind('.//ns:clrScheme', ns))
            color_tags = [c.tag.split('}')[-1] for c in clr_scheme]

            soup = bs4.BeautifulSoup(theme_data, 'xml')

            s_clr_scheme = soup.find('a:clrScheme')
            colors = [c for c in s_clr_scheme]
            for color in soup.find('a:clrScheme'):
                c_name = color.name
                c_spec = next(color.children)
                spec_name = c_spec.name
                spec_attr, spec_val = next(((k, v) for k, v in c_spec.attrs.items()))

            colors_with_srgb_spec = (color for color in soup.find('a:clrScheme'))
            colors_with_srgb_spec = [c for c in soup.find('a:clrScheme') if next(c.children).name == 'srgbClr']
            color_specs = {c.name: next(c.children).attrs['val'] for c in colors_with_srgb_spec}

            print(soup)

    parser = argparse.ArgumentParser()
    parser.add_argument('presentation', help='presentation to work on', completer=_pptx_completer)

    @cmd2.with_category('PPT')
    @cmd2.with_argparser(parser)
    def do_rpr_info(self, arg: argparse.Namespace):
        """
        Analyze rPr properties in PPT XML
        """
        ppt_path = arg.presentation
        if not os.path.isfile(ppt_path):
            self.perror(f'File "{ppt_path}" not found.')
            return
        if os.path.splitext(ppt_path)[1].lower() != '.pptx':
            self.perror(f'File "{ppt_path}" is not a PPTX file.')
            return

        def rpr_info(slide: str, content: bytes):
            """

            :param content:
            :return:
            """
            root = etree.fromstring(content)
            tree = etree.ElementTree(root)
            nsmap = root.nsmap
            rpr_infos = root.xpath('//a:rPr', namespaces=nsmap)
            for rpr_info in rpr_infos:
                path = tree.getpath(rpr_info)
                parent = rpr_info.getparent()
                solid_fill = rpr_info.find('a:solidFill', namespaces=nsmap)
                if solid_fill is not None:
                    color_spec = solid_fill[0]
                    if color_spec.tag.endswith('schemeClr'):
                        color_text = f'scheme: {color_spec.attrib["val"]}'
                    elif color_spec.tag.endswith('srgbClr'):
                        color_text = f'RGB: {color_spec.attrib["val"]}'
                    else:
                        raise KeyError
                    fill_text = color_text
                else:
                    fill_text = '<no fill>'
                # let's look at the parent
                t = rpr_info.xpath('../a:t', namespaces=nsmap)
                if t:
                    t_text = t[0].text
                else:
                    t_text = 'No text?'
                print(f'{slide}: {path} {fill_text} text: {t_text}')

        with ZipFile(file=ppt_path) as ppt_file:
            file_list = ppt_file.namelist()
            slides = [n for n in file_list if n.startswith('ppt/slides/slide')]
            for slide in slides:
                print()
                print(slide)
                with ppt_file.open(name=slide) as slide_file:
                    slide_content = slide_file.read()
                    rpr_info(slide=slide, content=slide_content)
        return

    parser = argparse.ArgumentParser()
    parser.add_argument('presentation', help='presentation to work on', completer=_pptx_completer)

    @cmd2.with_category('PPT')
    @cmd2.with_argparser(parser)
    def do_rels(self, arg: argparse.Namespace):
        """
        Analyze rels of PPTX
        """
        ppt_path = arg.presentation
        if not os.path.isfile(ppt_path):
            self.perror(f'File "{ppt_path}" not found.')
            return
        if os.path.splitext(ppt_path)[1].lower() != '.pptx':
            self.perror(f'File "{ppt_path}" is not a PPTX file.')
            return

        self.poutput('Slide dependencies:')
        with PPTXHelper(file=ppt_path) as ppt_file:
            for s_name, slide in ppt_file.slides.items():
                slide_layout = ppt_file.slide_layouts[slide.rel.slide_layout]
                slide_master = slide_layout.rel.slide_master
                self.poutput(f'{s_name}: slide layout {slide.rel.slide_layout} master {slide_master}')

            self.poutput('\n' * 3)
            p_rel = ppt_file.presentation_rel
            self.poutput(f'Presentation rel: {p_rel}')
            c_spec = ppt_file.color_spec
            self.poutput()
            self.poutput(f' Color spec: {c_spec}')

            for slide in p_rel.slide:
                slide_name = f'ppt/{slide}'
                slide_rel = ppt_file.slide_rel(slide_name=slide_name)
                self.poutput(f'{slide}: {slide_rel}')


if __name__ == '__main__':
    app = PPTApp()
    sys.exit(app.cmdloop())
